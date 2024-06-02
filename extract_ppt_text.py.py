from pptx import Presentation

# Carregar a apresentação
ppt = Presentation('C:\\Users\\marin\\Downloads\\A10_PENSAMENTO_SISTEMICO_NAS_ORGANIZAÇÕES_E_CIBERNÁTICA.pptx')

# Abrir um novo arquivo de texto para escrita
with open('C:\\Users\\marin\\Downloads\\apresentacao_extraida.txt', 'w', encoding='utf-8') as file:
    # Iterar pelos slides e pelos elementos de texto
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Escrever o texto no arquivo
                file.write(shape.text + '\n')

print("Texto extraído e salvo em apresentacao_extraida.txt")

